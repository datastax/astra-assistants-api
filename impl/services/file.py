import csv
import mimetypes
import os
import traceback
from io import BufferedReader
from typing import Optional, List

import docx2txt
import pptx
from fastapi import UploadFile, HTTPException
from loguru import logger
from PyPDF2 import PdfReader
from markitdown import MarkItDown, UnsupportedFormatException
from openai import OpenAI

from impl.astra_vector import HandledResponse
from impl.models import Document

exclude_exts: List[str] = [
    ".map",
    ".tfstate",
    ".jar",
    ".png",
    ".odg",
    ".bz2",
    ".xz",
    ".fits",
    ".jpg",
    ".jpeg",
    ".download",
    ".gif",
    ".bmp",
    ".tiff",
    ".ico",
    ".mp3",
    ".wav",
    ".wma",
    ".ogg",
    ".flac",
    ".mp4",
    ".avi",
    ".mkv",
    ".mov",
    ".patch",
    ".wmv",
    ".m4a",
    ".m4v",
    ".3gp",
    ".3g2",
    ".rm",
    ".swf",
    ".flv",
    ".iso",
    ".bin",
    ".tar",
    ".zip",
    ".7z",
    ".gz",
    ".rar",
    ".svg",
    ".pyc",
    ".pub",
    ".pem",
    ".ttf",
    ".dfn",
    ".dfm",
    ".feature",
    ".lock",
]

async def get_document_from_file(file: UploadFile, file_id: str, openai_api_key: str = None) -> Document:
    extracted_text = await extract_text_from_from_file(file, openai_api_key)

    doc = Document(id=file_id, text=extracted_text)

    return doc


def extract_text_from_filepath(filepath: str, mimetype: Optional[str] = None, openai_api_key: str = None) -> str:
    """Return the text content of a file given its filepath."""
    if mimetype is None or mimetype == "application/octet-stream":
        # Get the mimetype of the file based on its extension
        logger.info(f"guessing mimetype from filepath: {filepath}")
        mimetype, _ = mimetypes.guess_type(filepath.strip())
        logger.info(f"guessed: {mimetype}")
    #get extension from filepath for example /tmp/pytest.ini, safe because splittext returns a tuple
    extension = os.path.splitext(filepath)[1]
    if not mimetype:
        # when there's no mimetype, treat other valid extensions as text/plain, including files without extensions (i.e. Dockerfile)
        if extension not in exclude_exts:
            mimetype = "text/plain"
        else:
            # Unsupported file type
            raise HTTPException(
                status_code=400,
                detail="Unsupported file type: {}".format(filepath),
            )
    else:
        # treat programming language extensions as text/plain regardless of mimetype
        if extension in (".c", ".cpp", ".css", ".html", ".java", ".js", ".json", ".md", ".php", ".py", ".rb", ".ts", ".xml"):
            mimetype = "text/plain"
    try:
        with open(filepath, "rb") as file:
            extracted_text = extract_text_from_file(file, mimetype, openai_api_key, filepath)
    except Exception as e:
        logger.error(e)
        raise e

    return extracted_text


def extract_text_from_file(file: BufferedReader, mimetype: str, openai_api_key, filepath: str) -> str:
    filetype = mimetype
    logger.info(f"mimetype: {mimetype}")
    logger.info(f"filepath: {filepath}")
    if mimetype == "application/pdf":
        # Extract text from pdf using PyPDF2
        reader = PdfReader(file)
        extracted_text = " ".join([page.extract_text() for page in reader.pages])
    elif mimetype == "text/plain" or mimetype == "text/markdown" or mimetype == "application/sql" or mimetype == "message/rfc822":
        # Read text from plain text file
        extracted_text = file.read().decode("utf-8")
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        # Extract text from docx using docx2txt
        extracted_text = docx2txt.process(file)
    # TODO: supported formats should be Supported formats: "c", "cpp", "css", "csv", "docx", "gif", "html", "java", "jpeg", "jpg", "js", "json", "md", "pdf", "php", "png", "pptx", "py", "rb", "tar", "tex", "ts", "txt", "xlsx", "xml", "zip"
    # figure out what they do with the images.
    elif mimetype == "text/csv":
        # Extract text from csv using csv module
        extracted_text = ""
        decoded_buffer = (line.decode("utf-8") for line in file)
        reader = csv.reader(decoded_buffer)
        for row in reader:
            extracted_text += " ".join(row) + "\n"
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        # Extract text from pptx using python-pptx
        extracted_text = ""
        presentation = pptx.Presentation(file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            extracted_text += run.text + " "
                    extracted_text += "\n"
    else:
        if openai_api_key is None:
            raise HTTPException(
                status_code=400,
                detail="OpenAI API key is required for this file type: {}\nRemember to set your OPENAI_API_KEY env var".format(mimetype),
            )
        client = OpenAI(api_key=openai_api_key)
        md = MarkItDown(llm_client=client, llm_model="gpt-4o")
        try:
            extracted_text = md.convert(source=filepath).text_content
        except UnsupportedFormatException as e:
            logger.error(f"markitdown unsupported exception for: {mimetype}\nerror: {e}")
            raise HTTPException(
                status_code=400,
                detail="Unsupported file type: {}\nError: {}".format(mimetype, e)
            )
        except Exception as e:
            logger.error(f"Error: {e}\nTrace: {traceback.format_exc()}")
            # Unsupported file type
            raise HTTPException(
                status_code=400,
                detail="Unknown error parsing file type: {}\nError: {}".format(mimetype, e)
            )
    return extracted_text


# Extract text from a file based on its mimetype
async def extract_text_from_from_file(file: UploadFile, openai_api_key: str = None) -> str:
    """Return the text content of a file."""
    # get the file body from the upload file object
    mimetype = file.content_type
    logger.info(f"mimetype: {mimetype}")
    logger.info(f"file.file: {file.file}")
    logger.info(f"file: {file}")

    file_stream = await file.read()

    temp_file_path = "/tmp/"+file.filename

    # write the file to a temporary location
    with open(temp_file_path, "wb") as f:
        f.write(file_stream)

    try:
        extracted_text = extract_text_from_filepath(temp_file_path, mimetype, openai_api_key)
    except Exception as e:
        logger.error(e)
        os.remove(temp_file_path)
        raise e

    # remove file from temp location
    os.remove(temp_file_path)

    return extracted_text
